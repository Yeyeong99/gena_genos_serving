"""Office 문서 추출/주입/저장 런타임."""

from __future__ import annotations

from datetime import date, datetime
import re
from typing import Any, Dict, List, Optional, Tuple

from lxml import etree
from openpyxl import load_workbook
from openpyxl.utils.datetime import from_excel
from openpyxl.utils import get_column_letter, quote_sheetname
from openpyxl.utils.cell import range_boundaries
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn
from pptx.util import Pt

from translation_pipeline.common.nodes import PREVIEW_HEIGHT, PREVIEW_WIDTH, is_translatable

# 번역된 PPTX run 의 East Asian / Complex Script / Sym typeface 를 강제로
# 한글-호환 폰트(Noto Sans CJK KR) 로 덮어쓴다. 원본 PPTX 가 <a:ea typeface="Arial"/>
# 같은 라틴 폰트를 그대로 갖고 있으면 LibreOffice 가 한글 텍스트에 Liberation Sans
# 로 substitute → CJK 글리프 부재 → PDF/SVG 미리보기에서 default glyph("시" 류)
# 가 반복 출력되는 회귀가 발생한다 (#TODO post-mortem 참조). latin 은 라틴 문자에만
# 적용되므로 원본 디자인 유지를 위해 건드리지 않는다.
_KOREAN_PPTX_TYPEFACE = "Noto Sans CJK KR"
_KOREAN_TYPEFACE_TAGS = ("a:ea", "a:cs", "a:sym")


def _force_korean_typeface_on_run(run: Any, typeface: str = _KOREAN_PPTX_TYPEFACE) -> None:
    """번역된 PPTX run 의 ea/cs/sym typeface 를 한글-호환 폰트로 강제한다."""

    r = getattr(run, "_r", None)
    if r is None:
        return
    rPr = r.find(qn("a:rPr"))
    if rPr is None:
        rPr = etree.SubElement(r, qn("a:rPr"))
        r.insert(0, rPr)
    for tag in _KOREAN_TYPEFACE_TAGS:
        el = rPr.find(qn(tag))
        if el is None:
            el = etree.SubElement(rPr, qn(tag))
        el.set("typeface", typeface)


def _clamp(value: float, min_value: float, max_value: float) -> float:
    return max(min_value, min(max_value, value))


def _preview_bbox(x: float, y: float, w: float, h: float) -> List[int]:
    x0 = _clamp(x, 0, PREVIEW_WIDTH - 8)
    y0 = _clamp(y, 0, PREVIEW_HEIGHT - 8)
    x1 = _clamp(x0 + max(24, w), x0 + 8, PREVIEW_WIDTH)
    y1 = _clamp(y0 + max(18, h), y0 + 8, PREVIEW_HEIGHT)
    return [int(round(x0)), int(round(y0)), int(round(x1)), int(round(y1))]


def _safe_float(value: Any) -> Optional[float]:
    try:
        if value is None:
            return None
        return float(value)
    except (TypeError, ValueError):
        return None


def _normalize_text(value: Any) -> str:
    if value is None:
        return ""
    return str(value)


def _pptx_shape_bbox(shape: Any, presentation: Any) -> List[int]:
    slide_w = float(getattr(presentation, "slide_width", 1) or 1)
    slide_h = float(getattr(presentation, "slide_height", 1) or 1)
    left = float(getattr(shape, "left", 0) or 0)
    top = float(getattr(shape, "top", 0) or 0)
    width = float(getattr(shape, "width", 0) or 0)
    height = float(getattr(shape, "height", 0) or 0)
    return _preview_bbox(
        left / slide_w * PREVIEW_WIDTH,
        top / slide_h * PREVIEW_HEIGHT,
        width / slide_w * PREVIEW_WIDTH,
        height / slide_h * PREVIEW_HEIGHT,
    )


def _xlsx_column_units(sheet: Any, col_index: int) -> float:
    width = getattr(sheet.column_dimensions.get(get_column_letter(col_index)), "width", None)
    if width is None:
        width = 8.43
    return max(1.0, float(width))


def _xlsx_row_units(sheet: Any, row_index: int) -> float:
    height = getattr(sheet.row_dimensions.get(row_index), "height", None)
    if height is None:
        height = 15.0
    return max(1.0, float(height))


def _xlsx_merge_bounds(sheet: Any) -> Dict[Tuple[int, int], Tuple[int, int, int, int]]:
    merged: Dict[Tuple[int, int], Tuple[int, int, int, int]] = {}
    for merged_range in sheet.merged_cells.ranges:
        min_col, min_row, max_col, max_row = range_boundaries(str(merged_range))
        bounds = (min_row, min_col, max_row, max_col)
        for row in range(min_row, max_row + 1):
            for col in range(min_col, max_col + 1):
                merged[(row, col)] = bounds
    return merged


_KOREAN_WEEKDAYS_LONG = ["월요일", "화요일", "수요일", "목요일", "금요일", "토요일", "일요일"]
_KOREAN_WEEKDAYS_SHORT = ["월", "화", "수", "목", "금", "토", "일"]
_ENGLISH_WEEKDAYS_LONG = ["Monday", "Tuesday", "Wednesday", "Thursday", "Friday", "Saturday", "Sunday"]
_ENGLISH_WEEKDAYS_SHORT = ["Mon", "Tue", "Wed", "Thu", "Fri", "Sat", "Sun"]
_HANGUL_FINANCIAL_DIGITS = {
    0: "",
    1: "壹",
    2: "貳",
    3: "參",
    4: "四",
    5: "伍",
    6: "六",
    7: "七",
    8: "八",
    9: "九",
}
_HANGUL_FINANCIAL_SMALL_UNITS = ["", "拾", "百", "阡"]
_HANGUL_FINANCIAL_BIG_UNITS = ["", "萬", "億", "兆"]


def _xlsx_format_section(number_format: str, value: Any, *, text_value: bool = False) -> str:
    sections = (number_format or "").split(";")
    if text_value and len(sections) >= 4:
        return sections[3]
    if not sections:
        return ""
    try:
        numeric_value = float(value)
    except (TypeError, ValueError):
        return sections[0]
    if numeric_value > 0:
        return sections[0]
    if numeric_value < 0 and len(sections) >= 2:
        return sections[1]
    if numeric_value == 0 and len(sections) >= 3:
        return sections[2]
    return sections[0]


def _xlsx_hanja_amount(value: int) -> str:
    if value == 0:
        return "零"

    def render_group(group_value: int) -> str:
        parts: list[str] = []
        digits = [
            group_value // 1000 % 10,
            group_value // 100 % 10,
            group_value // 10 % 10,
            group_value % 10,
        ]
        units = ["阡", "百", "拾", ""]
        for digit, unit in zip(digits, units):
            if digit == 0:
                continue
            if digit == 1 and unit in {"拾", "百", "阡"}:
                parts.append(unit)
            else:
                parts.append(f"{_HANGUL_FINANCIAL_DIGITS[digit]}{unit}")
        return "".join(parts)

    parts: list[str] = []
    group_index = 0
    remaining = abs(value)
    while remaining:
        group = remaining % 10000
        if group:
            rendered = render_group(group)
            big_unit = _HANGUL_FINANCIAL_BIG_UNITS[group_index] if group_index < len(_HANGUL_FINANCIAL_BIG_UNITS) else ""
            parts.insert(0, f"{rendered}{big_unit}")
        remaining //= 10000
        group_index += 1
    return "".join(parts)


def _xlsx_format_arabic_number(value: int | float, placeholder: str) -> str:
    try:
        number = float(value)
    except (TypeError, ValueError):
        return str(value)
    decimal_places = 0
    if "." in placeholder:
        decimal_part = placeholder.rsplit(".", 1)[1]
        decimal_places = sum(1 for char in decimal_part if char in "0#?")
    use_comma = "," in placeholder
    if decimal_places > 0:
        formatted = f"{number:,.{decimal_places}f}" if use_comma else f"{number:.{decimal_places}f}"
        if "#" in placeholder.rsplit(".", 1)[1]:
            formatted = formatted.rstrip("0").rstrip(".")
    else:
        rounded = int(round(number))
        formatted = f"{rounded:,}" if use_comma else str(rounded)
    if formatted.startswith("0.") and placeholder.startswith("#"):
        formatted = formatted[1:]
    return formatted


def _xlsx_render_format_pattern(
    value: Any,
    number_format: str,
    *,
    text_value: bool = False,
) -> str:
    fmt = _xlsx_format_section(number_format, value, text_value=text_value)
    if not fmt or fmt.lower() == "general":
        return str(value)

    dbnum_match = re.search(r"\[dbnum(?P<kind>\d+)\]", fmt, flags=re.IGNORECASE)
    dbnum_kind = int(dbnum_match.group("kind")) if dbnum_match else None
    placeholder_inserted = False
    recognized = False
    rendered: list[str] = []
    index = 0

    def render_placeholder(placeholder: str) -> str:
        nonlocal placeholder_inserted, recognized
        placeholder_inserted = True
        recognized = True
        if dbnum_kind == 2:
            try:
                return _xlsx_hanja_amount(int(round(float(value))))
            except (TypeError, ValueError):
                return str(value)
        return _xlsx_format_arabic_number(value, placeholder)

    while index < len(fmt):
        char = fmt[index]
        if char == '"':
            end = fmt.find('"', index + 1)
            if end == -1:
                rendered.append(fmt[index + 1 :])
                break
            rendered.append(fmt[index + 1 : end])
            index = end + 1
            continue
        if char == "[":
            end = fmt.find("]", index + 1)
            index = len(fmt) if end == -1 else end + 1
            continue
        if char == "\\":
            if index + 1 < len(fmt):
                rendered.append(fmt[index + 1])
                index += 2
            else:
                index += 1
            continue
        if char in {"_", "*"}:
            index += 2
            continue
        if text_value and char == "@":
            rendered.append(str(value))
            placeholder_inserted = True
            recognized = True
            index += 1
            continue
        if not text_value and char in "0#?":
            end = index + 1
            while end < len(fmt) and fmt[end] in "0#?,.":
                end += 1
            rendered.append(render_placeholder(fmt[index:end]))
            index = end
            continue
        rendered.append(char)
        index += 1

    if text_value and not placeholder_inserted:
        return str(value)
    if not text_value and dbnum_kind and not placeholder_inserted:
        rendered.append(render_placeholder("0"))
    return "".join(rendered).strip() if recognized or dbnum_kind else str(value)


def _xlsx_number_format_has_text_literal(number_format: str) -> bool:
    fmt = _xlsx_format_section(number_format, 1)
    if re.search(r"\[dbnum\d+\]", fmt, flags=re.IGNORECASE):
        return True
    without_brackets = re.sub(r"\[[^\]]*\]", "", fmt)
    without_quotes = re.sub(r'"[^"]*"', "", without_brackets)
    literals = re.sub(r"[0#?,._*\\/@Ee+\-\s;:$€£¥₩%()]", "", without_quotes)
    quoted_literals = "".join(re.findall(r'"([^"]*)"', fmt))
    return bool(re.search(r"[^\d\s.,+\-/%()]", literals + quoted_literals))


def _xlsx_text_format_has_placeholder(number_format: str) -> bool:
    return "@" in _xlsx_format_section(number_format, "", text_value=True)


def _xlsx_format_date_text(value: date | datetime, number_format: str) -> str:
    fmt = (number_format or "").split(";", 1)[0]
    month = value.month
    day = value.day
    year = value.year
    weekday = value.weekday()
    long_weekday = _KOREAN_WEEKDAYS_LONG[weekday]
    short_weekday = _KOREAN_WEEKDAYS_SHORT[weekday]
    long_weekday_en = _ENGLISH_WEEKDAYS_LONG[weekday]
    short_weekday_en = _ENGLISH_WEEKDAYS_SHORT[weekday]

    def render_token(token: str) -> str:
        lower = token.lower()
        size = len(token)
        if lower.startswith("y"):
            if size <= 2:
                return f"{year % 100:02d}"
            return f"{year:04d}"
        if lower.startswith("m"):
            if size == 1:
                return str(month)
            if size == 2:
                return f"{month:02d}"
            if size in (3, 4):
                # Korean Excel commonly displays mmm/mmmm as "12월".
                return f"{month}월"
            return str(month)[0]
        if lower.startswith("d"):
            if size == 1:
                return str(day)
            if size == 2:
                return f"{day:02d}"
            if size == 3:
                return short_weekday_en
            return long_weekday_en
        if lower.startswith("a"):
            if size >= 4:
                return long_weekday
            return short_weekday
        return token

    rendered: list[str] = []
    index = 0
    recognized = False
    while index < len(fmt):
        char = fmt[index]
        if char == '"':
            end = fmt.find('"', index + 1)
            if end == -1:
                rendered.append(fmt[index + 1 :])
                break
            rendered.append(fmt[index + 1 : end])
            index = end + 1
            continue
        if char == "[":
            end = fmt.find("]", index + 1)
            index = len(fmt) if end == -1 else end + 1
            continue
        if char == "\\":
            if index + 1 < len(fmt):
                rendered.append(fmt[index + 1])
                index += 2
            else:
                index += 1
            continue
        if char == "_":
            index += 2
            continue
        if char == "*":
            index += 2
            continue
        if char.isalpha():
            end = index + 1
            while end < len(fmt) and fmt[end].lower() == char.lower():
                end += 1
            token = fmt[index:end]
            if token.lower()[0] in {"y", "m", "d", "a"}:
                rendered.append(render_token(token))
                recognized = True
            else:
                rendered.append(token)
            index = end
            continue
        rendered.append(char)
        index += 1

    if recognized:
        return "".join(rendered).strip()

    if "aaa" in fmt and "일" in fmt and "(" in fmt:
        return f"{day}일({short_weekday})"
    if "aaaa" in fmt and "일" not in fmt and "월" not in fmt:
        return long_weekday
    if "aaa" in fmt and "일" not in fmt and "월" not in fmt:
        return short_weekday
    if "월" in fmt and "일" in fmt:
        month_text = f"{month:02d}" if "mm" in fmt.lower() else str(month)
        day_text = f"{day:02d}" if "dd" in fmt.lower() else str(day)
        return f"{month_text}월 {day_text}일"
    if "월" in fmt:
        month_text = f"{month:02d}" if "mm" in fmt.lower() else str(month)
        return f"{month_text}월"
    if "일" in fmt:
        day_text = f"{day:02d}" if "dd" in fmt.lower() else str(day)
        return f"{day_text}일"
    return value.isoformat()


def _xlsx_number_format_looks_like_date(number_format: str) -> bool:
    fmt = (number_format or "").split(";", 1)[0]
    if not fmt or fmt.lower() == "general":
        return False
    normalized = re.sub(r"\[[^\]]*\]", "", fmt)
    normalized = re.sub(r'"[^"]*"', "", normalized)
    lower = normalized.lower()
    return any(token in lower for token in ("y", "d", "aaa", "월", "일")) or bool(
        re.search(r"(^|[^a-z])m{1,5}([^a-z]|$)", lower)
    )


def _xlsx_display_text(cell: Any, cached_cell: Any | None = None) -> str:
    raw_value = cell.value
    cached_value = getattr(cached_cell, "value", None) if cached_cell is not None else None
    value = cached_value if isinstance(raw_value, str) and raw_value.startswith("=") else raw_value
    if value is None and cached_value is not None:
        value = cached_value
    if value is None:
        return ""
    number_format = str(getattr(cell, "number_format", "") or "")
    if isinstance(value, str):
        if _xlsx_text_format_has_placeholder(number_format):
            return _xlsx_render_format_pattern(value, number_format, text_value=True)
        return value
    if isinstance(value, (datetime, date)):
        return _xlsx_format_date_text(value, number_format)
    if getattr(cell, "is_date", False) and isinstance(cached_value, (datetime, date)):
        return _xlsx_format_date_text(cached_value, number_format)
    if isinstance(value, (int, float)) and _xlsx_number_format_looks_like_date(number_format):
        try:
            epoch = getattr(getattr(cell, "parent", None).parent, "epoch", None)
            converted = from_excel(value, epoch=epoch) if epoch is not None else from_excel(value)
            if isinstance(converted, (datetime, date)):
                return _xlsx_format_date_text(converted, number_format)
        except Exception:
            return ""
    if isinstance(value, (int, float)) and _xlsx_number_format_has_text_literal(number_format):
        return _xlsx_render_format_pattern(value, number_format)
    return ""


def _xlsx_cell_bbox(
    sheet: Any,
    row: int,
    col: int,
    row_span: int = 1,
    col_span: int = 1,
) -> List[int]:
    base_cell_w = 108.0
    base_cell_h = 32.0
    default_col_units = 8.43
    default_row_units = 15.0

    x_units = sum(_xlsx_column_units(sheet, index) for index in range(1, max(1, col)))
    y_units = sum(_xlsx_row_units(sheet, index) for index in range(1, max(1, row)))
    w_units = sum(_xlsx_column_units(sheet, index) for index in range(col, col + max(1, col_span)))
    h_units = sum(_xlsx_row_units(sheet, index) for index in range(row, row + max(1, row_span)))

    x = 62 + x_units / default_col_units * base_cell_w
    y = 58 + y_units / default_row_units * base_cell_h
    width = w_units / default_col_units * base_cell_w
    height = h_units / default_row_units * base_cell_h
    return _preview_bbox(x, y, width, height)


def extract_docx(file_path: str) -> Tuple[Any, List[dict]]:
    """DOCX 문서에서 번역 가능한 텍스트 노드를 추출한다.

    Args:
        file_path: 입력 DOCX 경로.

    Returns:
        저장 컨텍스트와 노드 목록.
    """

    import zipfile

    from lxml import etree

    ns_w = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"
    ns_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
    ns_mc = "http://schemas.openxmlformats.org/markup-compatibility/2006"

    def is_inside_fallback(elem: Any) -> bool:
        parent = elem.getparent()
        while parent is not None:
            if parent.tag == f"{{{ns_mc}}}Fallback":
                return True
            parent = parent.getparent()
        return False

    def collect_t_from_runs(container: Any) -> List[Any]:
        t_nodes: List[Any] = []
        for child in container:
            tag = child.tag
            if tag == f"{{{ns_w}}}r":
                for sub in child:
                    if sub.tag == f"{{{ns_w}}}t":
                        t_nodes.append(sub)
            elif tag == f"{{{ns_w}}}hyperlink":
                for wr in child:
                    if wr.tag == f"{{{ns_w}}}r":
                        for sub in wr:
                            if sub.tag == f"{{{ns_w}}}t":
                                t_nodes.append(sub)
            elif tag == f"{{{ns_w}}}ins":
                for wr in child:
                    if wr.tag == f"{{{ns_w}}}r":
                        for sub in wr:
                            if sub.tag == f"{{{ns_w}}}t":
                                t_nodes.append(sub)
        return t_nodes

    def extract_from_tree(root: Any, source: str) -> List[dict]:
        found: List[dict] = []
        for wp in root.iter(f"{{{ns_w}}}p"):
            if is_inside_fallback(wp):
                continue
            t_nodes = collect_t_from_runs(wp)
            text = "".join(t.text for t in t_nodes if t.text).strip()
            if text and is_translatable(text):
                found.append({"type": "xml_text", "t_nodes": t_nodes, "text": text, "source": source})
            for sdt in wp:
                if sdt.tag != f"{{{ns_w}}}sdt":
                    continue
                if is_inside_fallback(sdt):
                    continue
                sdt_content = sdt.find(f"{{{ns_w}}}sdtContent")
                if sdt_content is None:
                    continue
                sdt_t = collect_t_from_runs(sdt_content)
                sdt_text = "".join(t.text for t in sdt_t if t.text).strip()
                if sdt_text and is_translatable(sdt_text):
                    found.append({"type": "xml_text", "t_nodes": sdt_t, "text": sdt_text, "source": source})

        for ap in root.iter(f"{{{ns_a}}}p"):
            if is_inside_fallback(ap):
                continue
            t_nodes = [t for t in ap.iter(f"{{{ns_a}}}t") if not is_inside_fallback(t)]
            text = "".join(t.text for t in t_nodes if t.text).strip()
            if text and is_translatable(text):
                found.append({"type": "xml_text", "t_nodes": t_nodes, "text": text, "source": source})
        return found

    nodes: List[dict] = []
    xml_parts: Dict[str, Any] = {}
    parser = etree.XMLParser(remove_blank_text=False)

    with zipfile.ZipFile(file_path, "r") as zip_file:
        doc_root = etree.fromstring(zip_file.read("word/document.xml"), parser)
        xml_parts["word/document.xml"] = doc_root
        nodes.extend(extract_from_tree(doc_root, "body"))

        for name in zip_file.namelist():
            if (name.startswith("word/header") or name.startswith("word/footer")) and name.endswith(".xml"):
                hf_root = etree.fromstring(zip_file.read(name), parser)
                xml_parts[name] = hf_root
                source = "header" if "header" in name else "footer"
                nodes.extend(extract_from_tree(hf_root, source))

    context = {"file_path": file_path, "xml_parts": xml_parts}
    print(f"[DOCX 추출] {len(nodes)}개 텍스트 추출 완료")
    return context, nodes


def extract_xlsx(file_path: str) -> Tuple[Any, List[dict]]:
    """XLSX 문서에서 번역 가능한 셀 노드를 추출한다.

    Args:
        file_path: 입력 XLSX 경로.

    Returns:
        워크북 객체와 노드 목록.
    """

    workbook = load_workbook(file_path, data_only=False)
    setattr(workbook, "_ai_translation_source_path", file_path)
    cached_workbook = load_workbook(file_path, data_only=True)
    cached_sheets = {sheet.title: sheet for sheet in cached_workbook.worksheets}
    nodes: List[dict] = []
    for sheet in workbook.worksheets:
        cached_sheet = cached_sheets.get(sheet.title)
        merged_bounds = _xlsx_merge_bounds(sheet)
        for row in sheet.iter_rows():
            for cell in row:
                raw_value = cell.value
                if isinstance(raw_value, str) and raw_value.startswith("="):
                    continue
                cached_cell = cached_sheet[cell.coordinate] if cached_sheet is not None else None
                value = _xlsx_display_text(cell, cached_cell)
                if is_translatable(value):
                    min_row, min_col, max_row, max_col = merged_bounds.get(
                        (cell.row, cell.column),
                        (cell.row, cell.column, cell.row, cell.column),
                    )
                    row_span = max_row - min_row + 1
                    col_span = max_col - min_col + 1
                    nodes.append(
                        {
                            "type": "cell",
                            "group": "sheet_cell",
                            "cell": cell,
                            "text": value,
                            "sheet_name": sheet.title,
                            "row": cell.row,
                            "col": cell.column,
                            "cell_ref": cell.coordinate,
                            "row_span": row_span,
                            "col_span": col_span,
                            "merged_range": (
                                f"{get_column_letter(min_col)}{min_row}:{get_column_letter(max_col)}{max_row}"
                                if row_span > 1 or col_span > 1
                                else ""
                            ),
                            "bbox": _xlsx_cell_bbox(sheet, min_row, min_col, row_span, col_span),
                        }
                    )
    print(f"[XLSX 추출] {len(nodes)}개 번역 노드 추출 완료")
    return workbook, nodes


def extract_pptx(file_path: str) -> Tuple[Any, List[dict]]:
    """PPTX 문서에서 번역 가능한 텍스트 노드를 추출한다.

    Args:
        file_path: 입력 PPTX 경로.

    Returns:
        프레젠테이션 객체와 노드 목록.
    """

    presentation = Presentation(file_path)
    nodes: List[dict] = []

    def process_text_frame(
        text_frame: Any,
        slide_index: int,
        shape: Any,
        group: str,
        bbox: Optional[List[int]] = None,
    ) -> None:
        shape_name = _normalize_text(getattr(shape, "name", ""))
        shape_bbox = bbox or _pptx_shape_bbox(shape, presentation)
        for paragraph in text_frame.paragraphs:
            full_text = "".join(run.text for run in paragraph.runs)
            if full_text.strip() and is_translatable(full_text):
                nodes.append(
                    {
                        "type": "paragraph",
                        "group": group,
                        "paragraph": paragraph,
                        "text": full_text,
                        "slide_index": slide_index,
                        "shape_name": shape_name,
                        "bbox": shape_bbox,
                    }
                )

    def process_shape(shape: Any, slide_index: int) -> None:
        shape_name = _normalize_text(getattr(shape, "name", ""))
        try:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                for child in shape.shapes:
                    process_shape(child, slide_index)
            elif shape.has_table:
                for row_index, row in enumerate(shape.table.rows):
                    for col_index, cell in enumerate(row.cells):
                        if getattr(cell, "is_spanned", False):
                            continue
                        if cell.text_frame:
                            try:
                                slide_w = float(getattr(presentation, "slide_width", 1) or 1)
                                slide_h = float(getattr(presentation, "slide_height", 1) or 1)
                                col_span = int(getattr(cell, "span_width", 1) or 1)
                                row_span = int(getattr(cell, "span_height", 1) or 1)
                                x_offset = sum(float(shape.table.columns[index].width) for index in range(col_index))
                                y_offset = sum(float(shape.table.rows[index].height) for index in range(row_index))
                                cell_w = sum(
                                    float(shape.table.columns[index].width)
                                    for index in range(col_index, min(len(shape.table.columns), col_index + col_span))
                                )
                                cell_h = sum(
                                    float(shape.table.rows[index].height)
                                    for index in range(row_index, min(len(shape.table.rows), row_index + row_span))
                                )
                                cell_bbox = _preview_bbox(
                                    (float(shape.left) + x_offset) / slide_w * PREVIEW_WIDTH,
                                    (float(shape.top) + y_offset) / slide_h * PREVIEW_HEIGHT,
                                    cell_w / slide_w * PREVIEW_WIDTH,
                                    cell_h / slide_h * PREVIEW_HEIGHT,
                                )
                            except Exception:
                                cell_bbox = _pptx_shape_bbox(shape, presentation)
                            process_text_frame(cell.text_frame, slide_index, shape, "table_cell", cell_bbox)
            elif shape.has_chart:
                chart = shape.chart
                chart_bbox = _pptx_shape_bbox(shape, presentation)
                if chart.has_title and chart.chart_title.has_text_frame:
                    process_text_frame(chart.chart_title.text_frame, slide_index, shape, "chart_title", chart_bbox)
                try:
                    if chart.value_axis.has_title:
                        process_text_frame(chart.value_axis.axis_title.text_frame, slide_index, shape, "chart_axis", chart_bbox)
                    if chart.category_axis.has_title:
                        process_text_frame(
                            chart.category_axis.axis_title.text_frame,
                            slide_index,
                            shape,
                            "chart_axis",
                            chart_bbox,
                        )
                except (ValueError, AttributeError):
                    pass
                try:
                    plots = chart.plots
                    if plots:
                        for index, category in enumerate(list(plots[0].categories)):
                            if isinstance(category, str) and is_translatable(category):
                                nodes.append(
                                    {
                                        "type": "chart_category",
                                        "group": "chart_category",
                                        "chart": chart,
                                        "index": index,
                                        "text": category,
                                        "slide_index": slide_index,
                                        "shape_name": shape_name,
                                        "bbox": chart_bbox,
                                    }
                                )
                        for index, series in enumerate(chart.series):
                            if isinstance(series.name, str) and is_translatable(series.name):
                                nodes.append(
                                    {
                                        "type": "chart_series",
                                        "group": "chart_series",
                                        "chart": chart,
                                        "index": index,
                                        "text": series.name,
                                        "slide_index": slide_index,
                                        "shape_name": shape_name,
                                        "bbox": chart_bbox,
                                    }
                                )
                except Exception:
                    pass
            elif shape.has_text_frame:
                process_text_frame(shape.text_frame, slide_index, shape, "text_frame")
        except Exception as exc:
            print(f"    [경고] Shape 처리 중 오류: {exc}")

    for slide_index, slide in enumerate(presentation.slides, start=1):
        for shape in slide.shapes:
            process_shape(shape, slide_index)

    print(f"[PPTX 추출] {len(nodes)}개 노드 추출 완료")
    return presentation, nodes


def _node_translation(node: dict, trans_map: Dict[str, str]) -> str:
    """노드에 적용할 번역 문자열을 구한다.

    Args:
        node: 번역 대상 노드.
        trans_map: 원문/번역 매핑.

    Returns:
        주입할 번역 문자열.
    """

    if node.get("translated_text") is not None:
        return str(node.get("translated_text", ""))
    original = str(node.get("text", ""))
    return str(trans_map.get(original, original))


def _run_style_key(run: Any) -> tuple[Any, Any, Any, Any]:
    font = getattr(run, "font", None)
    return (
        getattr(font, "size", None),
        getattr(font, "bold", None),
        getattr(font, "italic", None),
        getattr(font, "name", None),
    )


def _group_runs_by_style(runs: Any) -> List[List[int]]:
    groups: List[List[int]] = []
    previous_key: tuple[Any, Any, Any, Any] | None = None
    for index, run in enumerate(runs):
        key = _run_style_key(run)
        if not groups or key != previous_key:
            groups.append([index])
        else:
            groups[-1].append(index)
        previous_key = key
    return groups


def _split_title_and_body_translation(translated: str) -> tuple[str, str] | None:
    # A PPT paragraph can contain a bold title and smaller bullet body in
    # separate runs. The LLM often fuses them as "... documents: - Load ...".
    # Split at the first bullet marker so the original run styles survive.
    match = re.search(
        r"(?P<prefix>[:：]?\s*)(?P<body>(?:[-–—―－‒]\s*|•\s+))",
        translated,
    )
    if match:
        head = translated[: match.start()].rstrip()
        body = translated[match.start("body") :].strip()
    else:
        colon_match = re.search(r"[:：]\s+(?=[A-Z0-9])", translated)
        if not colon_match:
            return None
        head = translated[: colon_match.start()].rstrip()
        body = translated[colon_match.end() :].strip()
    if head.endswith((':', '：')):
        head = head[:-1].rstrip()
    if not head or not body:
        return None
    return head, body


def _apply_text_to_run_group(runs: Any, group: List[int], text: str) -> None:
    if not group:
        return
    runs[group[0]].text = text
    if text:
        _force_korean_typeface_on_run(runs[group[0]])
    for index in group[1:]:
        runs[index].text = ""


def _apply_translated_paragraph_text(paragraph: Any, translated: str, node: dict) -> None:
    runs = paragraph.runs
    if not runs:
        return

    groups = _group_runs_by_style(runs)
    split = _split_title_and_body_translation(translated) if len(groups) >= 2 else None
    if split:
        head, body = split
        original_body_text = "".join(runs[index].text for index in groups[1]).strip()
        if original_body_text.startswith(("-", "–", "•")) and not body.startswith(("-", "–", "•")):
            body = f"- {body}"
        _apply_text_to_run_group(runs, groups[0], head)
        _apply_text_to_run_group(runs, groups[1], body)
        for group in groups[2:]:
            _apply_text_to_run_group(runs, group, "")
    else:
        _apply_text_to_run_group(runs, list(range(len(runs))), translated)

    if node.get("edited_font_size") is not None:
        try:
            runs[0].font.size = Pt(float(node["edited_font_size"]))
        except (TypeError, ValueError):
            pass


def inject_docx(context: Any, nodes: List[dict], trans_map: Dict[str, str]) -> None:
    """DOCX XML 텍스트 노드에 번역 결과를 주입한다.

    Args:
        context: DOCX 저장 컨텍스트.
        nodes: 번역 노드 목록.
        trans_map: 원문/번역 매핑.

    Returns:
        없음.
    """

    count = 0
    for node in nodes:
        original = str(node["text"])
        translated = _node_translation(node, trans_map)
        if not translated or translated == original:
            continue
        t_nodes = node["t_nodes"]
        if t_nodes:
            t_nodes[0].text = translated
            for item in t_nodes[1:]:
                item.text = ""
            count += 1
    print(f"[DOCX 주입] {count}개 노드 번역 적용")


def save_docx(context: Any, output_path: str) -> None:
    """DOCX 수정 XML을 다시 파일로 저장한다.

    Args:
        context: DOCX 저장 컨텍스트.
        output_path: 저장할 경로.

    Returns:
        없음.
    """

    import zipfile

    from lxml import etree

    file_path = context["file_path"]
    xml_parts = context["xml_parts"]
    with zipfile.ZipFile(file_path, "r") as source_zip:
        with zipfile.ZipFile(output_path, "w", zipfile.ZIP_DEFLATED) as target_zip:
            for item in source_zip.infolist():
                if item.filename in xml_parts:
                    xml_bytes = etree.tostring(
                        xml_parts[item.filename],
                        xml_declaration=True,
                        encoding="UTF-8",
                        standalone=True,
                    )
                    target_zip.writestr(item, xml_bytes)
                else:
                    target_zip.writestr(item, source_zip.read(item.filename))


def inject_xlsx(workbook: Any, nodes: List[dict], trans_map: Dict[str, str]) -> None:
    """XLSX 셀에 번역 결과를 주입한다.

    Args:
        workbook: 워크북 객체.
        nodes: 셀 노드 목록.
        trans_map: 원문/번역 매핑.

    Returns:
        없음.
    """

    count = 0
    sheet_rename_count = 0
    sheet_rename_map: Dict[str, str] = {}
    existing_sheet_names = set(getattr(workbook, "sheetnames", []) or [])

    # 1. 시트 이름을 먼저 확정한다. 이후 셀 수식 참조와 셀 번역은 확정된 시트명 기준으로 처리한다.
    for node in nodes:
        if node.get("type") != "sheet_name":
            continue
        original = str(node["text"])
        translated = _node_translation(node, trans_map)
        if not translated or translated == original:
            continue

        sheet = node.get("sheet")
        if sheet is None:
            continue
        original_title = str(getattr(sheet, "title", "") or node.get("sheet_name", ""))
        new_title = _unique_xlsx_sheet_title(translated, original_title, existing_sheet_names)
        if new_title and new_title != original_title:
            existing_sheet_names.discard(original_title)
            sheet.title = new_title
            existing_sheet_names.add(new_title)
            sheet_rename_map[original_title] = new_title
            sheet_rename_count += 1

    # 2. 시트명 참조가 들어간 수식/이름 정의를 먼저 갱신한다.
    if sheet_rename_map:
        _rewrite_xlsx_sheet_references(workbook, sheet_rename_map)

    # 3. 수식이 아닌 일반 셀 번역을 적용한다.
    for node in nodes:
        if node.get("type") == "sheet_name":
            continue
        original = str(node["text"])
        translated = _node_translation(node, trans_map)
        if not translated or translated == original:
            continue

        cell = node["cell"]
        cell.value = translated
        count += 1
    print(f"[XLSX 주입] {count}개 셀 번역 적용, {sheet_rename_count}개 시트명 변경")


def _unique_xlsx_sheet_title(
    title: str,
    current_title: str,
    existing_sheet_names: set[str],
) -> str:
    """Excel sheet title 제약을 맞추고 중복을 피한다."""

    cleaned = re.sub(r"[\[\]:*?/\\]", " ", title).strip() or current_title
    cleaned = re.sub(r"\s+", " ", cleaned)
    base = cleaned[:31] or current_title[:31]
    if base == current_title or base not in existing_sheet_names:
        return base

    for index in range(2, 100):
        suffix = f" ({index})"
        candidate = f"{base[:31 - len(suffix)]}{suffix}"
        if candidate not in existing_sheet_names:
            return candidate
    return base[:28] + "..."


def _rewrite_xlsx_sheet_references(workbook: Any, rename_map: Dict[str, str]) -> None:
    """시트명 변경 후 수식/이름 정의 안의 시트 참조를 함께 갱신한다."""

    def replace_reference_text(value: str) -> str:
        updated = value
        for old_name, new_name in rename_map.items():
            old_quoted = quote_sheetname(old_name)
            new_quoted = quote_sheetname(new_name)
            updated = updated.replace(f"{old_quoted}!", f"{new_quoted}!")
            updated = re.sub(
                rf"(?<![\w\]']){re.escape(old_name)}(?=!)",
                new_quoted,
                updated,
            )
        return updated

    for worksheet in workbook.worksheets:
        for row in worksheet.iter_rows():
            for cell in row:
                value = cell.value
                if isinstance(value, str) and value.startswith("="):
                    cell.value = replace_reference_text(value)

        data_validations = getattr(getattr(worksheet, "data_validations", None), "dataValidation", None)
        if data_validations:
            for validation in data_validations:
                for attr in ("formula1", "formula2"):
                    value = getattr(validation, attr, None)
                    if isinstance(value, str):
                        setattr(validation, attr, replace_reference_text(value))

    defined_names = getattr(workbook, "defined_names", None)
    values_fn = getattr(defined_names, "values", None)
    if callable(values_fn):
        for defined_name in values_fn():
            text = getattr(defined_name, "attr_text", None)
            if isinstance(text, str):
                defined_name.attr_text = replace_reference_text(text)


def inject_pptx(presentation: Any, nodes: List[dict], trans_map: Dict[str, str]) -> None:
    """PPTX 텍스트 노드와 차트 데이터에 번역 결과를 주입한다.

    Args:
        presentation: 프레젠테이션 객체.
        nodes: 노드 목록.
        trans_map: 원문/번역 매핑.

    Returns:
        없음.
    """

    _ = presentation
    count = 0
    chart_updates: Dict[int, dict] = {}

    for node in nodes:
        node_type = node["type"]
        original = str(node["text"])
        translated = _node_translation(node, trans_map)
        if not translated or translated == original:
            continue

        if node_type == "paragraph":
            paragraph = node["paragraph"]
            if paragraph.runs:
                _apply_translated_paragraph_text(paragraph, translated, node)
                count += 1
        elif node_type == "chart_category":
            chart = node["chart"]
            chart_id = id(chart)
            chart_updates.setdefault(chart_id, {"chart": chart, "categories": {}, "series": {}})
            chart_updates[chart_id]["categories"][node["index"]] = translated
        elif node_type == "chart_series":
            chart = node["chart"]
            chart_id = id(chart)
            chart_updates.setdefault(chart_id, {"chart": chart, "categories": {}, "series": {}})
            chart_updates[chart_id]["series"][node["index"]] = translated

    for info in chart_updates.values():
        chart = info["chart"]
        try:
            plots = chart.plots
            if not plots:
                continue
            original_categories = list(plots[0].categories)
            translated_categories = [
                info["categories"].get(index, category)
                for index, category in enumerate(original_categories)
            ]

            chart_data = CategoryChartData()
            chart_data.categories = translated_categories
            for index, series in enumerate(chart.series):
                chart_data.add_series(info["series"].get(index, series.name), series.values)
            chart.replace_data(chart_data)
            count += len(info["categories"]) + len(info["series"])
        except Exception as exc:
            print(f"    [차트 데이터 교체 실패] {exc}")

    print(f"[PPTX 주입] {count}개 노드 번역 적용")
