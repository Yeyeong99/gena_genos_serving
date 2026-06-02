"""PPTX 추출/주입 런타임."""

from __future__ import annotations

from typing import Any, Dict, List, Optional, Tuple

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER

from translation_pipeline.common.logging_utils import log_info
from translation_pipeline.common.nodes import PREVIEW_HEIGHT, PREVIEW_WIDTH, is_translatable

from .pptx_style import _apply_translated_paragraph_text
from .runtime_common import _element_type_with_placeholder, _node_translation, _preview_bbox


def _pptx_shape_bbox(shape: Any, presentation: Any) -> List[int]:
    slide_w = float(getattr(presentation, "slide_width", 1) or 1)
    slide_h = float(getattr(presentation, "slide_height", 1) or 1)
    left = float(getattr(shape, "left", 0) or 0)
    top = float(getattr(shape, "top", 0) or 0)
    width = float(getattr(shape, "width", 0) or 0)
    height = float(getattr(shape, "height", 0) or 0)
    return _preview_bbox(
        left / slide_w * PREVIEW_WIDTH,
        top / slide_h * PREVIEW_HEIGHT,
        width / slide_w * PREVIEW_WIDTH,
        height / slide_h * PREVIEW_HEIGHT,
    )



def extract_pptx(file_path: str) -> Tuple[Any, List[dict]]:
    """PPTX 문서에서 번역 가능한 텍스트 노드를 추출한다.

    Args:
        file_path: 입력 PPTX 경로.

    Returns:
        프레젠테이션 객체와 노드 목록.
    """

    presentation = Presentation(file_path)
    nodes: List[dict] = []
    table_index = 0

    def pptx_shape_element_type(shape: Any, slide: Any) -> str:
        if getattr(slide.shapes, "title", None) is shape:
            return "slide_title"
        if getattr(shape, "is_placeholder", False):
            try:
                placeholder_type = shape.placeholder_format.type
                title_types = {
                    getattr(PP_PLACEHOLDER, "TITLE", None),
                    getattr(PP_PLACEHOLDER, "CENTER_TITLE", None),
                }
                body_types = {
                    getattr(PP_PLACEHOLDER, "BODY", None),
                    getattr(PP_PLACEHOLDER, "CONTENT", None),
                    getattr(PP_PLACEHOLDER, "OBJECT", None),
                }
                if placeholder_type in title_types:
                    return "slide_title"
                if placeholder_type in body_types:
                    return "text_box"
            except Exception:
                pass
        return "text_box"

    def process_text_frame(
        text_frame: Any,
        slide_index: int,
        shape: Any,
        group: str,
        bbox: Optional[List[int]] = None,
        metadata: Optional[dict[str, Any]] = None,
    ) -> None:
        shape_name = str(getattr(shape, "name", "") or "")
        shape_bbox = bbox or _pptx_shape_bbox(shape, presentation)
        metadata = metadata or {}
        for paragraph in text_frame.paragraphs:
            full_text = "".join(run.text for run in paragraph.runs)
            if full_text.strip() and is_translatable(full_text):
                element_type = _element_type_with_placeholder(
                    full_text,
                    str(metadata.get("element_type") or group),
                )
                nodes.append(
                    {
                        "type": "paragraph",
                        "group": group,
                        "paragraph": paragraph,
                        "text": full_text,
                        "slide_index": slide_index,
                        "shape_name": shape_name,
                        "bbox": shape_bbox,
                        "doc_format": "pptx",
                        **metadata,
                        "element_type": element_type,
                    }
                )

    def process_shape(shape: Any, slide_index: int, slide: Any) -> None:
        nonlocal table_index
        shape_name = str(getattr(shape, "name", "") or "")
        try:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                for child in shape.shapes:
                    process_shape(child, slide_index, slide)
            elif shape.has_table:
                current_table_index = table_index
                table_index += 1
                for row_index, row in enumerate(shape.table.rows):
                    for col_index, cell in enumerate(row.cells):
                        if getattr(cell, "is_spanned", False):
                            continue
                        if cell.text_frame:
                            try:
                                slide_w = float(getattr(presentation, "slide_width", 1) or 1)
                                slide_h = float(getattr(presentation, "slide_height", 1) or 1)
                                col_span = int(getattr(cell, "span_width", 1) or 1)
                                row_span = int(getattr(cell, "span_height", 1) or 1)
                                x_offset = sum(float(shape.table.columns[index].width) for index in range(col_index))
                                y_offset = sum(float(shape.table.rows[index].height) for index in range(row_index))
                                cell_w = sum(
                                    float(shape.table.columns[index].width)
                                    for index in range(col_index, min(len(shape.table.columns), col_index + col_span))
                                )
                                cell_h = sum(
                                    float(shape.table.rows[index].height)
                                    for index in range(row_index, min(len(shape.table.rows), row_index + row_span))
                                )
                                cell_bbox = _preview_bbox(
                                    (float(shape.left) + x_offset) / slide_w * PREVIEW_WIDTH,
                                    (float(shape.top) + y_offset) / slide_h * PREVIEW_HEIGHT,
                                    cell_w / slide_w * PREVIEW_WIDTH,
                                    cell_h / slide_h * PREVIEW_HEIGHT,
                                )
                            except Exception:
                                cell_bbox = _pptx_shape_bbox(shape, presentation)
                            process_text_frame(
                                cell.text_frame,
                                slide_index,
                                shape,
                                "table_cell",
                                cell_bbox,
                                {
                                    "table_index": current_table_index,
                                    "row_index": row_index,
                                    "col_index": col_index,
                                    "row": row_index,
                                    "col": col_index,
                                    "is_header": row_index == 0,
                                },
                            )
            elif shape.has_chart:
                chart = shape.chart
                chart_bbox = _pptx_shape_bbox(shape, presentation)
                if chart.has_title and chart.chart_title.has_text_frame:
                    process_text_frame(
                        chart.chart_title.text_frame,
                        slide_index,
                        shape,
                        "chart_title",
                        chart_bbox,
                        {"element_type": "chart_title"},
                    )
                try:
                    if chart.value_axis.has_title:
                        process_text_frame(
                            chart.value_axis.axis_title.text_frame,
                            slide_index,
                            shape,
                            "chart_axis",
                            chart_bbox,
                            {"element_type": "chart_axis"},
                        )
                    if chart.category_axis.has_title:
                        process_text_frame(
                            chart.category_axis.axis_title.text_frame,
                            slide_index,
                            shape,
                            "chart_axis",
                            chart_bbox,
                            {"element_type": "chart_axis"},
                        )
                except (ValueError, AttributeError):
                    pass
                try:
                    plots = chart.plots
                    if plots:
                        for index, category in enumerate(list(plots[0].categories)):
                            if isinstance(category, str) and is_translatable(category):
                                nodes.append(
                                    {
                                        "type": "chart_category",
                                        "group": "chart_category",
                                        "chart": chart,
                                        "index": index,
                                        "text": category,
                                        "slide_index": slide_index,
                                        "shape_name": shape_name,
                                        "bbox": chart_bbox,
                                        "doc_format": "pptx",
                                        "element_type": "chart_category",
                                    }
                                )
                        for index, series in enumerate(chart.series):
                            if isinstance(series.name, str) and is_translatable(series.name):
                                nodes.append(
                                    {
                                        "type": "chart_series",
                                        "group": "chart_series",
                                        "chart": chart,
                                        "index": index,
                                        "text": series.name,
                                        "slide_index": slide_index,
                                        "shape_name": shape_name,
                                        "bbox": chart_bbox,
                                        "doc_format": "pptx",
                                        "element_type": "chart_series",
                                    }
                                )
                except Exception:
                    pass
            elif shape.has_text_frame:
                process_text_frame(
                    shape.text_frame,
                    slide_index,
                    shape,
                    "text_frame",
                    metadata={"element_type": pptx_shape_element_type(shape, slide)},
                )
        except Exception as exc:
            log_info(f"    [경고] Shape 처리 중 오류: {exc}")

    for slide_index, slide in enumerate(presentation.slides, start=1):
        for shape in slide.shapes:
            process_shape(shape, slide_index, slide)

    log_info(f"[PPTX 추출] {len(nodes)}개 노드 추출 완료")
    return presentation, nodes



def inject_pptx(presentation: Any, nodes: List[dict], trans_map: Dict[str, str]) -> None:
    """PPTX 텍스트 노드와 차트 데이터에 번역 결과를 주입한다.

    Args:
        presentation: 프레젠테이션 객체.
        nodes: 노드 목록.
        trans_map: 원문/번역 매핑.

    Returns:
        없음.
    """

    _ = presentation
    count = 0
    chart_updates: Dict[int, dict] = {}

    for node in nodes:
        node_type = node["type"]
        original = str(node["text"])
        translated = _node_translation(node, trans_map)
        if not translated or translated == original:
            continue

        if node_type == "paragraph":
            paragraph = node["paragraph"]
            if paragraph.runs:
                _apply_translated_paragraph_text(paragraph, translated, node)
                count += 1
        elif node_type == "chart_category":
            chart = node["chart"]
            chart_id = id(chart)
            chart_updates.setdefault(chart_id, {"chart": chart, "categories": {}, "series": {}})
            chart_updates[chart_id]["categories"][node["index"]] = translated
        elif node_type == "chart_series":
            chart = node["chart"]
            chart_id = id(chart)
            chart_updates.setdefault(chart_id, {"chart": chart, "categories": {}, "series": {}})
            chart_updates[chart_id]["series"][node["index"]] = translated

    for info in chart_updates.values():
        chart = info["chart"]
        try:
            plots = chart.plots
            if not plots:
                continue
            original_categories = list(plots[0].categories)
            translated_categories = [
                info["categories"].get(index, category)
                for index, category in enumerate(original_categories)
            ]

            chart_data = CategoryChartData()
            chart_data.categories = translated_categories
            for index, series in enumerate(chart.series):
                chart_data.add_series(info["series"].get(index, series.name), series.values)
            chart.replace_data(chart_data)
            count += len(info["categories"]) + len(info["series"])
        except Exception as exc:
            log_info(f"    [차트 데이터 교체 실패] {exc}")

    log_info(f"[PPTX 주입] {count}개 노드 번역 적용")
